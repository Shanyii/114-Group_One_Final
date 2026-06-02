"""
@module document_parser
@description 文件解析工具。支援 PDF（pdfplumber）與 PPT/PPTX（python-pptx）格式，
             解析後依 config.chunk_size 切分為 Chunk，供 RAG 向量化使用。
@dependencies pdfplumber, python-pptx, core.config
@author 楊沁霖（RAG / 資料處理）
@version 1.1.0
"""

import logging
import os
import re
import uuid
from pathlib import Path

import aiosqlite
import pdfplumber
from pptx import Presentation

from core.config import get_settings

logger = logging.getLogger(__name__)


class RecursiveCharacterTextSplitter:
    """
    智慧型遞迴字元切分器。
    優先以段落 (\\n\\n)、換行 (\\n)、中文標點 (。、？、！、；、，、、) 等邊界進行切分，
    確保切分出的每一個 Chunk 長度都不超過 chunk_size，且盡可能保留語意完整性。
    """

    def __init__(self, chunk_size: int, chunk_overlap: int) -> None:
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        # 依重要性優先級排序的分隔符號
        self.separators = [
            "\n\n", "\n", "。", "？", "！", "；", "；", "，", "、",
            ". ", "? ", "! ", ", ", ",", " ", ""
        ]

    def split_text(self, text: str) -> list[str]:
        """遞迴切分入口。"""
        return self._split_text(text, self.separators)

    def _split_text(self, text: str, separators: list[str]) -> list[str]:
        # 如果文字長度已經小於等於 chunk_size，直接返回
        if len(text) <= self.chunk_size:
            return [text]

        # 如果沒有分隔符號了，強行按 chunk_size 切片
        if not separators:
            chunks = []
            start = 0
            while start < len(text):
                chunks.append(text[start : start + self.chunk_size])
                start += self.chunk_size - self.chunk_overlap
            return chunks

        # 選擇當前層級的分隔符號
        separator = separators[0]
        next_separators = separators[1:]

        # 使用分隔符號切割
        if separator == "":
            splits = list(text)
        else:
            splits = text.split(separator)

        leaf_chunks = []
        for i, split in enumerate(splits):
            # 保留分隔符號（如標點符號不該被丟失，應加回前一個句子的尾部）
            punctuation_separators = ["。", "？", "！", "；", "，", "、", ". ", "? ", "! ", ", ", ","]
            if i < len(splits) - 1 and separator in punctuation_separators:
                split_content = split + separator
            else:
                split_content = split

            if not split_content.strip():
                continue

            # 遞迴切分
            leaf_chunks.extend(self._split_text(split_content, next_separators))

        # 將葉子節點合併為滿足 chunk_size 和 chunk_overlap 的 chunks
        merged_chunks = []
        current_pieces = []
        current_len = 0

        for piece in leaf_chunks:
            piece_len = len(piece)
            if piece_len > self.chunk_size:
                if current_pieces:
                    merged_chunks.append("".join(current_pieces).strip())
                    current_pieces = []
                    current_len = 0
                merged_chunks.append(piece)
                continue

            if current_len + piece_len > self.chunk_size:
                if current_pieces:
                    merged_chunks.append("".join(current_pieces).strip())

                # 計算 overlap：保留部分舊的 pieces
                # 限制：除了長度不超過 chunk_overlap 外，新 Chunk (overlap + piece) 的總長度也絕不能超過 chunk_size
                overlap_pieces = []
                overlap_len = 0
                for p in reversed(current_pieces):
                    if (overlap_len + len(p) <= self.chunk_overlap) and (overlap_len + len(p) + piece_len <= self.chunk_size):
                        overlap_pieces.insert(0, p)
                        overlap_len += len(p)
                    else:
                        break

                current_pieces = overlap_pieces + [piece]
                current_len = overlap_len + piece_len
            else:
                current_pieces.append(piece)
                current_len += piece_len

        if current_pieces:
            merged_chunks.append("".join(current_pieces).strip())

        return merged_chunks

# ── 支援的 MIME Type 白名單 ───────────────────────────────────────────────────
ALLOWED_MIME_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
}

MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50 MB


class DocumentParser:
    """
    文件解析工具。

    負責：
    1. 驗證並儲存上傳檔案（至 uploads/ 暫存目錄）
    2. 解析 PDF / PPT 為純文字
    3. 切分文字為 Chunk 列表（用於 Embedding）
    4. 將解析結果存入 SQLite documents 資料表

    Example:
        >>> parser = DocumentParser()
        >>> doc_id = await parser.save_and_parse(file_bytes, "lecture.pdf", student_id)
        >>> chunks = await parser.get_chunks(doc_id)
    """

    def __init__(self) -> None:
        self.settings = get_settings()

    async def save_and_parse(
        self,
        file_bytes: bytes,
        filename: str,
        student_id: str,
        content_type: str,
    ) -> str:
        """
        儲存上傳檔案並解析為純文字，存入資料庫。

        Args:
            file_bytes: 上傳檔案的二進位內容
            filename: 原始檔案名稱
            student_id: 學生 UUID
            content_type: MIME Type

        Returns:
            str: 生成的 document_id（UUID）

        Raises:
            ValueError: 檔案格式不支援或超過大小限制
        """
        # 驗證檔案大小
        if len(file_bytes) > MAX_FILE_SIZE_BYTES:
            raise ValueError(f"檔案大小超過限制（最大 50MB），目前：{len(file_bytes) / 1024 / 1024:.1f}MB")

        # 驗證 MIME Type
        file_type = ALLOWED_MIME_TYPES.get(content_type)
        if file_type is None:
            raise ValueError(f"不支援的檔案格式：{content_type}。僅接受 PDF 與 PPTX")

        document_id = str(uuid.uuid4())

        # 暫存至 uploads 目錄
        temp_path = os.path.join(self.settings.upload_dir, f"{document_id}.{file_type}")
        os.makedirs(self.settings.upload_dir, exist_ok=True)
        with open(temp_path, "wb") as f:
            f.write(file_bytes)

        # 解析純文字
        try:
            raw_text = self._parse_file(temp_path, file_type)
        finally:
            # 解析完成後刪除暫存檔（不長期儲存原始講義）
            if os.path.exists(temp_path):
                os.remove(temp_path)

        if not raw_text.strip():
            raise ValueError("無法從檔案中提取文字內容，請確認檔案未加密或非空白")

        # 存入資料庫
        await self._save_to_db(document_id, student_id, filename, file_type, raw_text)
        logger.info("[DocParser] 文件解析完成：doc_id=%s, 字元數=%d", document_id, len(raw_text))
        return document_id

    def _parse_file(self, file_path: str, file_type: str) -> str:
        """
        依檔案類型解析純文字。

        Args:
            file_path: 暫存檔案路徑
            file_type: 'pdf' | 'pptx'

        Returns:
            str: 提取的純文字
        """
        if file_type == "pdf":
            return self._parse_pdf(file_path)
        elif file_type == "pptx":
            return self._parse_pptx(file_path)
        raise ValueError(f"未知的檔案類型：{file_type}")

    def _parse_pdf(self, file_path: str) -> str:
        """使用 pdfplumber 提取 PDF 文字（支援中文）。"""
        texts = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text(x_tolerance=3, y_tolerance=3)
                    if text:
                        texts.append(f"--- 第 {page_num} 頁 ---\n{text}")
        except Exception as exc:
            raise ValueError(f"PDF 解析失敗：{exc}") from exc
        return "\n\n".join(texts)

    def _parse_pptx(self, file_path: str) -> str:
        """使用 python-pptx 提取簡報文字。"""
        texts = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f"--- 第 {slide_num} 張投影片 ---\n" + "\n".join(slide_texts))
        except Exception as exc:
            raise ValueError(f"PPTX 解析失敗：{exc}") from exc
        return "\n\n".join(texts)

    def chunk_text(self, text: str) -> list[dict]:
        """
        將純文字切分為智慧型的 Chunk，並帶有頁碼/投影片 Metadata。

        步驟：
        1. 透過 Regex 解析 PDF / PPTX 的頁碼標記
        2. 將文字分成不同頁面的段落
        3. 對每一頁使用 RecursiveCharacterTextSplitter 進行智慧切分
        4. 打上正確的 page_num metadata

        Args:
            text: 原始純文字（可能含 page 標記）

        Returns:
            list[dict]: 每個項目包含：
                - text: 切分出的文字 Chunk
                - page_num: 所在的頁碼（整數，預設為 1）
        """
        chunk_size = self.settings.chunk_size
        overlap = self.settings.chunk_overlap
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)

        # 透過 Regex 尋找 --- 第 X 頁 --- 或 --- 第 X 張投影片 ---
        page_pattern = r"--- 第 (\d+) (頁|張投影片) ---"
        matches = list(re.finditer(page_pattern, text))

        segments = []
        if not matches:
            # 沒有任何頁面標記，將整個文字視為第 1 頁
            segments.append((text, 1))
        else:
            # 處理第一個標記之前的內容（若有）
            first_start = matches[0].start()
            if first_start > 0:
                pre_text = text[:first_start].strip()
                if pre_text:
                    segments.append((pre_text, 1))

            # 處理標記之間的內容
            for i in range(len(matches)):
                current_match = matches[i]
                page_num = int(current_match.group(1))
                segment_start = current_match.end()
                segment_end = matches[i + 1].start() if i < len(matches) - 1 else len(text)

                page_text = text[segment_start:segment_end].strip()
                if page_text:
                    segments.append((page_text, page_num))

        chunks = []
        for page_text, page_num in segments:
            # 智慧遞迴切分這一頁的文字
            page_chunks = splitter.split_text(page_text)
            for pc in page_chunks:
                if pc.strip():
                    chunks.append({
                        "text": pc,
                        "page_num": page_num
                    })

        logger.info("[DocParser] 智慧切分完成：共 %d 個 Chunks", len(chunks))
        return chunks

    async def read_from_db(self, document_id: str) -> str:
        """
        從資料庫讀取文件純文字。

        Args:
            document_id: 文件 UUID

        Returns:
            str: 純文字內容

        Raises:
            ValueError: 文件不存在
        """
        async with aiosqlite.connect(self.settings.database_url) as db:
            async with db.execute(
                "SELECT raw_text FROM documents WHERE document_id = ?", (document_id,)
            ) as cursor:
                row = await cursor.fetchone()
                if row is None or not row[0]:
                    raise ValueError(f"找不到文件或文字為空：{document_id}")
                return row[0]

    async def get_chunks(self, document_id: str) -> list[dict]:
        """
        取得文件的 Chunk 列表（從 DB 讀取後切分，包含 page_num）。

        Args:
            document_id: 文件 UUID

        Returns:
            list[dict]: 帶有 page_num metadata 的 Chunk 列表
        """
        raw_text = await self.read_from_db(document_id)
        return self.chunk_text(raw_text)

    async def _save_to_db(
        self,
        document_id: str,
        student_id: str,
        filename: str,
        file_type: str,
        raw_text: str,
    ) -> None:
        """將解析結果存入 documents 資料表。"""
        from datetime import datetime
        async with aiosqlite.connect(self.settings.database_url) as db:
            await db.execute(
                """
                INSERT INTO documents
                (document_id, student_id, filename, file_type, raw_text, uploaded_at)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (document_id, student_id, filename, file_type, raw_text,
                 datetime.utcnow().isoformat()),
            )
            await db.commit()
